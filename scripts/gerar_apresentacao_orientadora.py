"""Gera apresentacao PPTX para reuniao com a orientadora — Fase 1 + 2 artigos."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "Apresentacao_Orientadora_Fase1_Artigos.pptx"
ASSETS = Path(r"C:\Users\Ever\.cursor\projects\c-java-doutorado-ZonasCalor\assets")
SCREEN = ASSETS / "c__Users_Ever_AppData_Roaming_Cursor_User_workspaceStorage_0ccbb252b5dbcb598f6ddfb00c884692_images_image-00460f1e-fa64-4f28-ae38-e555a65566ca.png"
IMG_ARCH = ASSETS / "slide-arquitetura-fase1.png"
IMG_ART = ASSETS / "slide-dois-artigos.png"

# Paleta institucional (sem roxo genérico)
BG = RGBColor(0x0F, 0x1C, 0x24)
PANEL = RGBColor(0x16, 0x28, 0x32)
ACCENT = RGBColor(0x3D, 0x9B, 0x8F)
ACCENT2 = RGBColor(0xE8, 0xB8, 0x6D)
INK = RGBColor(0xE8, 0xF1, 0xF4)
MUTED = RGBColor(0x9B, 0xB0, 0xBA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DANGER = RGBColor(0xC5, 0x3D, 0x3D)

W = Inches(13.333)
H = Inches(7.5)


def _set_run(run, size: int, bold: bool = False, color: RGBColor = INK, font: str = "Calibri") -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _fill_solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_bg(slide) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    _fill_solid(shape, BG)
    # move to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def _bar(slide, y=Inches(0)) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, W, Inches(0.08))
    _fill_solid(bar, ACCENT)


def _textbox(slide, left, top, width, height, text: str, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size, bold, color)
    return tf


def _bullets(slide, left, top, width, height, items: list[str], size=16, color=INK) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        _set_run(run, size, False, color)


def _card(slide, left, top, width, height, title: str, body_lines: list[str]) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill_solid(card, PANEL)
    card.adjustments[0] = 0.08
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    _fill_solid(accent, ACCENT)
    _textbox(slide, left + Inches(0.25), top + Inches(0.15), width - Inches(0.4), Inches(0.4), title, 16, True, ACCENT)
    _bullets(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.35), height - Inches(0.7), body_lines, 13, INK)


def new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg(slide)
    return slide


def build() -> Path:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1. Capa
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.4), "UNESP — PPGEP | Câmpus de Bauru", 14, False, MUTED)
    _textbox(
        s,
        Inches(0.7),
        Inches(2.1),
        Inches(12),
        Inches(1.4),
        "Otimização de Escalas e Roteamento Territorial\nno Consultório na Rua",
        32,
        True,
        WHITE,
    )
    _textbox(
        s,
        Inches(0.7),
        Inches(3.7),
        Inches(12),
        Inches(0.8),
        "Fase 1 do SAD: fundação espacial, IA preditiva e estratégia dos dois artigos",
        20,
        False,
        ACCENT,
    )
    _textbox(s, Inches(0.7), Inches(5.5), Inches(6), Inches(0.4), "Doutorando: Ever Santoro", 16, False, INK)
    _textbox(s, Inches(0.7), Inches(5.95), Inches(10), Inches(0.4), "Apresentação para orientadora — demonstração da divisão da tese", 14, False, MUTED)

    # 2. Objetivo da reunião
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5), "Objetivo desta reunião", 28, True, WHITE)
    _bullets(
        s,
        Inches(0.7),
        Inches(1.2),
        Inches(11.5),
        Inches(5.5),
        [
            "Apresentar como a tese será dividida em fases tecnológicas e científicas.",
            "Demonstrar o protótipo da Fase 1 já em execução (mapa de calor com IA).",
            "Explicar detalhadamente o que a Inteligência Artificial faz neste momento.",
            "Propor a estratégia de publicação em dois artigos Qualis A (Scopus / WoS).",
            "Alinhar com a orientadora o escopo, fronteiras e cronograma de cada artigo.",
        ],
        18,
    )

    # 3. Problema
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5), "Problema de pesquisa", 28, True, WHITE)
    _card(
        s,
        Inches(0.5),
        Inches(1.1),
        Inches(6.0),
        Inches(5.5),
        "Contexto social e operacional",
        [
            "Crescimento da População em Situação de Rua (PSR) no Brasil.",
            "Consultório na Rua (CnR): cuidado itinerante na Atenção Primária (SUS).",
            "Equipes multidisciplinares com escassez crônica de recursos.",
            "Escalas e rotas ainda planejadas de forma majoritariamente manual.",
            "Risco de baixa cobertura e de não engajar os casos mais vulneráveis.",
        ],
    )
    _card(
        s,
        Inches(6.8),
        Inches(1.1),
        Inches(6.0),
        Inches(5.5),
        "Lacuna científica e tecnológica",
        [
            "Há literatura de Home Health Care Routing (HHCRSP) e Nurse Rostering.",
            "Há ML para predição em saúde, em geral hospitalar ou epidemiológica.",
            "Falta integração IA + Pesquisa Operacional no contexto CnR / SUS / SUAS.",
            "Ambiente urbano dinâmico e de alta vulnerabilidade é pouco modelado.",
            "Esta tese propõe um SAD híbrido para preencher essa lacuna.",
        ],
    )

    # 4. Divisão da tese
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5), "Como a tese se divide (visão para a orientadora)", 26, True, WHITE)
    phases = [
        ("Fase 1", "Fundação\nespacial + IA\n+ MVP visual", True),
        ("Fase 2", "Aprimorar\npredição\ncom dados reais", False),
        ("Fase 3", "Modelagem\nNRP + VRP\n(otimização)", False),
        ("Fase 4–5", "Integração\ndo SAD +\nvalidação", False),
    ]
    x0 = 0.6
    for i, (title, body, current) in enumerate(phases):
        left = Inches(x0 + i * 3.15)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.4), Inches(2.9), Inches(4.2))
        _fill_solid(box, PANEL)
        box.adjustments[0] = 0.1
        topbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.4), Inches(2.9), Inches(0.55))
        _fill_solid(topbar, ACCENT if current else RGBColor(0x2A, 0x45, 0x52))
        _textbox(s, left, Inches(1.48), Inches(2.9), Inches(0.4), title, 18, True, WHITE, PP_ALIGN.CENTER)
        _textbox(s, left + Inches(0.15), Inches(2.3), Inches(2.6), Inches(2.5), body, 16, False, INK, PP_ALIGN.CENTER)
        if current:
            _textbox(s, left, Inches(5.2), Inches(2.9), Inches(0.35), "EM ANDAMENTO", 12, True, ACCENT2, PP_ALIGN.CENTER)
    _textbox(
        s,
        Inches(0.6),
        Inches(6.7),
        Inches(12),
        Inches(0.4),
        "A Fase 1 é a base empírica e tecnológica dos dois artigos — sem ela, a otimização não tem demanda espacial confiável.",
        14,
        False,
        MUTED,
    )

    # 5. Importância da Fase 1
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5), "Por que a Fase 1 é decisiva", 28, True, WHITE)
    _bullets(
        s,
        Inches(0.7),
        Inches(1.1),
        Inches(12),
        Inches(5.8),
        [
            "Transforma o problema social em problema formalizável: pontos, densidades, clusters e restrições.",
            "Cria a camada de dados georreferenciados (PostGIS) exigida por qualquer motor de roteamento/escalas.",
            "Introduz IA preditiva cedo, evitando que a otimização opere só de forma reativa e estática.",
            "Gera artefato demonstrável (MVP visual) para diálogo com gestores, SMS/SUAS e comitê acadêmico.",
            "Produz o primeiro núcleo publicável (Artigo 1): método territorial + pipeline + validação visual.",
            "Reduz risco científico: prova viabilidade técnica antes do exame de qualificação e do Artigo 2 (PO).",
            "Permite evolução ética: dados sintéticos agora; calibração com CadÚnico/IPEA/e-SUS depois (LGPD).",
        ],
        17,
    )

    # 6. Arquitetura
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.45), "Arquitetura da Fase 1 (já implementada)", 26, True, WHITE)
    if IMG_ARCH.exists():
        s.shapes.add_picture(str(IMG_ARCH), Inches(0.5), Inches(0.9), width=Inches(12.3))
    _textbox(
        s,
        Inches(0.6),
        Inches(6.6),
        Inches(12),
        Inches(0.5),
        "Separação clara: Python = inteligência e ingestão | Java = API e visualização | PostGIS = verdade espacial.",
        14,
        False,
        MUTED,
    )

    # 7. O que a IA faz — detalhe
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.45), "O que a Inteligência Artificial faz neste momento", 24, True, WHITE)
    _textbox(s, Inches(0.6), Inches(0.75), Inches(12), Inches(0.35), "Script: modelo_preditivo_ia.py  ·  Tipo: Machine Learning não supervisionado (não é LLM/chatbot)", 13, False, ACCENT2)

    steps = [
        ("1. Mock territorial", "Gera ~500 pontos sintéticos em Bauru, concentrados em 4 epicentros (centro, terminal, sul, norte), com nível de vulnerabilidade 1–5."),
        ("2. K-Means", "Agrupa coordenadas (lon, lat) em 4 clusters. Descobre aglomerados naturais no território — as futuras “zonas de calor”."),
        ("3. KDE", "Estima densidade espacial contínua (Kernel Density). Pontos em áreas densas recebem score de concentração mais alto."),
        ("4. Score híbrido", "probabilidade_ia = 0,65×densidade_normalizada + 0,35×(vulnerabilidade/5). Resultado em [0,1] por paciente."),
        ("5. Persistência", "Grava pacientes (coordenada + score) e zonas (polígono, centroide, intensidade) no PostGIS."),
        ("6. Visualização", "O SAD Java lê /api/predicoes/heatmap e plota Leaflet.heat + marcadores dos centroides."),
    ]
    for i, (t, b) in enumerate(steps):
        col = i % 2
        row = i // 2
        left = Inches(0.5 + col * 6.4)
        top = Inches(1.25 + row * 1.85)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.15), Inches(1.7))
        _fill_solid(card, PANEL)
        card.adjustments[0] = 0.08
        _textbox(s, left + Inches(0.2), top + Inches(0.15), Inches(5.7), Inches(0.35), t, 15, True, ACCENT)
        _textbox(s, left + Inches(0.2), top + Inches(0.55), Inches(5.7), Inches(1.0), b, 13, False, INK)

    # 8. Demo screenshot
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.4), "Demonstração: MVP visual com zonas de calor preditas", 24, True, WHITE)
    if SCREEN.exists():
        # wide screenshot
        s.shapes.add_picture(str(SCREEN), Inches(0.45), Inches(0.7), width=Inches(12.4))
    _textbox(
        s,
        Inches(0.5),
        Inches(6.85),
        Inches(12.3),
        Inches(0.4),
        "Protótipo em execução: 500 pacientes · 4 zonas IA · heatmap com probabilidade_ia · Bauru-SP",
        13,
        False,
        MUTED,
        PP_ALIGN.CENTER,
    )

    # 9. Leitura do mapa
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5), "Como interpretar o mapa na demonstração", 26, True, WHITE)
    _card(
        s,
        Inches(0.5),
        Inches(1.1),
        Inches(6.0),
        Inches(5.4),
        "Elementos visuais",
        [
            "Manchas azul/verde: menor densidade ou vulnerabilidade predita.",
            "Amarelo/laranja/vermelho: maior probabilidade_ia (prioridade territorial).",
            "Círculos: centroides das 4 zonas encontradas pelo K-Means.",
            "Popup da zona: intensidade, nº de pacientes e método (KMEANS_KDE).",
            "Botão “Atualizar heatmap”: recarrega a API após novo ciclo Python.",
        ],
    )
    _card(
        s,
        Inches(6.8),
        Inches(1.1),
        Inches(6.0),
        Inches(5.4),
        "Valor para a gestão do CnR",
        [
            "Indica onde as equipes deveriam priorizar presença itinerante.",
            "Antecipa polos de vulnerabilidade em vez de só reagir a demandas.",
            "Serve de entrada futura para o motor de roteamento (matriz de custos).",
            "Facilita comunicação com SMS/SUAS com evidência espacial.",
            "Dados atuais são sintéticos; a lógica do método já é a definitiva.",
        ],
    )

    # 10. Estratégia dos 2 artigos
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.4), "Estratégia científica: dois artigos Qualis A", 26, True, WHITE)
    if IMG_ART.exists():
        s.shapes.add_picture(str(IMG_ART), Inches(1.2), Inches(0.85), width=Inches(10.8))
    _textbox(
        s,
        Inches(0.6),
        Inches(6.5),
        Inches(12),
        Inches(0.6),
        "Divisão intencional: Artigo 1 = contribuição em IA espacial / DSS preditivo  ·  Artigo 2 = contribuição em Pesquisa Operacional / otimização integrada.",
        14,
        False,
        MUTED,
    )

    # 11. Artigo 1 detalhado
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.4), "Artigo 1 — Predição territorial e zonas de calor no CnR", 24, True, WHITE)
    _textbox(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.35), "Âncora: Fases 1–2  ·  Ênfase: IA + geoinformação + SAD visual", 13, False, ACCENT2)

    _card(
        s,
        Inches(0.4),
        Inches(1.15),
        Inches(4.1),
        Inches(5.5),
        "Pergunta e contribuição",
        [
            "Como antecipar polos de vulnerabilidade da PSR no território urbano para o CnR?",
            "Método híbrido K-Means + KDE + score de vulnerabilidade.",
            "Pipeline reproduzível Python → PostGIS → heatmap.",
            "MVP de DSS para gestores municipais.",
            "Discussão ética/LGPD e uso de dados sintéticos calibráveis.",
        ],
    )
    _card(
        s,
        Inches(4.7),
        Inches(1.15),
        Inches(4.1),
        Inches(5.5),
        "Estrutura proposta do paper",
        [
            "Introdução: CnR, PSR, lacuna preditiva.",
            "Revisão: spatial ML em saúde + outreach.",
            "Método: dados, K-Means, KDE, score, arquitetura.",
            "Resultados: clusters, mapas, indicadores.",
            "Discussão: utilidade gerencial e limites.",
            "Conclusão + agenda para otimização (Art. 2).",
        ],
    )
    _card(
        s,
        Inches(9.0),
        Inches(1.15),
        Inches(3.9),
        Inches(5.5),
        "Alvos e entregáveis",
        [
            "Alvos: Healthcare Management; IJMI; BMC Health Serv.; Spatial journals.",
            "Também: tracks de DSS em saúde.",
            "Entregáveis: dataset sintético documentado, código, figuras de heatmap, métricas de cluster.",
            "Validação futura: CadÚnico/IPEA/e-SUS.",
        ],
    )

    # 12. Artigo 2 detalhado
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.4), "Artigo 2 — Otimização de escalas e roteamento territorial", 24, True, WHITE)
    _textbox(s, Inches(0.6), Inches(0.7), Inches(12), Inches(0.35), "Âncora: Fases 3–5  ·  Ênfase: Pesquisa Operacional + integração com a demanda predita", 13, False, ACCENT2)

    _card(
        s,
        Inches(0.4),
        Inches(1.15),
        Inches(4.1),
        Inches(5.5),
        "Pergunta e contribuição",
        [
            "Como alocar equipes multidisciplinares e rotas itinerantes maximizando cobertura sob restrições do SUS?",
            "Modelo híbrido NRP (escalas) + VRP com janelas.",
            "Demanda alimentada pelas zonas de calor (Art. 1).",
            "Comparação com planejamento manual.",
            "Heurísticas/metaheurísticas (ex.: Timefold).",
        ],
    )
    _card(
        s,
        Inches(4.7),
        Inches(1.15),
        Inches(4.1),
        Inches(5.5),
        "Estrutura proposta do paper",
        [
            "Introdução: logística do CnR e escassez.",
            "Revisão: HHCRSP, NRP, VRP em saúde.",
            "Formulação matemática e restrições SUS.",
            "Integração com matriz OSRM + demanda IA.",
            "Experimentos: cobertura, tempo, esgotamento.",
            "Conclusão e implicações de política pública.",
        ],
    )
    _card(
        s,
        Inches(9.0),
        Inches(1.15),
        Inches(3.9),
        Inches(5.5),
        "Alvos e entregáveis",
        [
            "Alvos: Computers & Operations Research; EJOR; Health Care Manag. Sci.",
            "Entregáveis: modelo, instâncias, resultados vs baseline manual.",
            "Indicadores: cobertura, deslocamento, turnos válidos, fairness da carga.",
        ],
    )

    # 13. Cruzamento fase x artigo
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5), "Cruzamento: etapas da tese × conteúdo dos artigos", 24, True, WHITE)

    headers = ["Etapa", "Entrega tecnológica", "Artigo 1", "Artigo 2"]
    rows = [
        ["Fase 1 (atual)", "PostGIS + ML + MVP mapa", "Núcleo metodológico", "Insumo de demanda"],
        ["Fase 2", "Dados reais / calibração", "Validação empírica", "Atualiza nós de demanda"],
        ["Fase 3", "NRP + VRP + solvers", "Discussão futura", "Núcleo metodológico"],
        ["Fase 4–5", "SAD integrado + simulação", "Caso de uso DSS", "Experimentos e baseline"],
        ["Qualificação", "Texto + protótipo", "Draft / submissão", "Formulação avançada"],
        ["Defesa", "Tese completa", "Publicado/aceito", "Publicado/aceito"],
    ]
    # simple table as cards
    col_w = [Inches(2.2), Inches(3.6), Inches(3.4), Inches(3.4)]
    x = Inches(0.5)
    y = Inches(1.0)
    for j, h in enumerate(headers):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_w[j], Inches(0.45))
        _fill_solid(cell, ACCENT)
        _textbox(s, x, y + Inches(0.05), col_w[j], Inches(0.35), h, 13, True, WHITE, PP_ALIGN.CENTER)
        x += col_w[j]
    for i, row in enumerate(rows):
        x = Inches(0.5)
        y = Inches(1.45 + i * 0.85)
        for j, val in enumerate(row):
            cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_w[j], Inches(0.8))
            _fill_solid(cell, PANEL if i % 2 == 0 else RGBColor(0x12, 0x22, 0x2C))
            _textbox(s, x + Inches(0.08), y + Inches(0.2), col_w[j] - Inches(0.1), Inches(0.5), val, 12, j == 0, INK, PP_ALIGN.CENTER)
            x += col_w[j]

    # 14. Cronograma artigos
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5), "Roteiro de criação dos dois artigos", 26, True, WHITE)
    _card(
        s,
        Inches(0.5),
        Inches(1.1),
        Inches(6.0),
        Inches(5.5),
        "Artigo 1 — sequência de produção",
        [
            "1) Congelar método K-Means+KDE e métricas.",
            "2) Gerar figuras padronizadas (heatmap, clusters).",
            "3) Escrever método + resultados com dados sintéticos.",
            "4) Incorporar calibração CadÚnico/IPEA quando possível.",
            "5) Submeter a periódico de gestão/informática em saúde.",
            "6) Usar revisões para fortalecer a tese e o Art. 2.",
        ],
    )
    _card(
        s,
        Inches(6.8),
        Inches(1.1),
        Inches(6.0),
        Inches(5.5),
        "Artigo 2 — sequência de produção",
        [
            "1) Formalizar NRP+VRP com restrições do MS/SUS.",
            "2) Ligar demanda às zonas do Artigo 1.",
            "3) Implementar solver (heurística/Timefold) + OSRM.",
            "4) Benchmark vs planejamento manual.",
            "5) Analisar cobertura, tempo e carga das equipes.",
            "6) Submeter a periódico de Operations Research / Health Care OR.",
        ],
    )

    # 15. Status atual
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5), "Status atual (o que já está pronto)", 26, True, WHITE)
    _bullets(
        s,
        Inches(0.7),
        Inches(1.2),
        Inches(12),
        Inches(5.5),
        [
            "Docker + PostgreSQL/PostGIS com schema espacial (pacientes e zonas_calor_preditas).",
            "Pipeline de IA em Python operacional (mock + K-Means + KDE + gravação no banco).",
            "Backend Spring Boot com API /api/predicoes/heatmap, /zonas e /resumo.",
            "Frontend Leaflet com heatmap e centroides — demonstrável agora em localhost:8081.",
            "README passo a passo e repositório GitHub (doutorado-fase1).",
            "Próximo alinhamento desejado: validar fronteira Artigo 1 × Artigo 2 e venues preferidas.",
        ],
        17,
    )

    # 16. Pedidos à orientadora
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5), "Pontos para alinhamento com a orientadora", 26, True, WHITE)
    _bullets(
        s,
        Inches(0.7),
        Inches(1.2),
        Inches(12),
        Inches(5.5),
        [
            "A divisão Artigo 1 (IA espacial/DSS) × Artigo 2 (PO/otimização) está adequada ao PPGEP?",
            "Preferência de periódicos para cada artigo (lista Qualis A / impacto).",
            "Grau de aprofundamento teórico exigido na Fase 1 antes da qualificação.",
            "Estratégia de dados reais: parceria municipal, CadÚnico agregados ou apenas sintéticos calibrados?",
            "Ordem de submissão: Artigo 1 antes da qualificação é prioridade?",
            "Necessidade de coautoria/rede internacional para Computers & OR / Healthcare Management.",
        ],
        17,
    )

    # 17. Encerramento
    s = new_slide(prs)
    _bar(s)
    _textbox(s, Inches(0.7), Inches(2.2), Inches(12), Inches(0.6), "Síntese", 18, False, ACCENT)
    _textbox(
        s,
        Inches(0.7),
        Inches(2.8),
        Inches(12),
        Inches(1.5),
        "A Fase 1 já materializa a fundação do SAD:\nIA preditiva + dados espaciais + mapa de decisão.\nSobre ela se constroem dois artigos complementares e a otimização das fases seguintes.",
        22,
        True,
        WHITE,
    )
    _textbox(s, Inches(0.7), Inches(5.5), Inches(12), Inches(0.4), "Obrigado — aberto a orientações e ajustes de escopo.", 16, False, MUTED)
    _textbox(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.4), "Ever Santoro  ·  UNESP/PPGEP  ·  Consultório na Rua", 14, False, ACCENT)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Gerado: {path}")
